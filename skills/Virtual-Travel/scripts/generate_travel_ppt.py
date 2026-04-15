#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
旅行日记 PPT 生成器
功能：下载景点图片、创建文件夹结构、生成 PowerPoint 旅行日记
用法：python generate_travel_ppt.py --city 北京 --output-dir "c:\\Users\\Terry\\Desktop\\Virtual-Travel-Skills-main"
"""

import os
import sys
import json
import argparse
import requests
from pathlib import Path
from datetime import datetime

try:
    from PIL import Image
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Cm
except ImportError:
    print("缺少依赖，请运行：pip install python-pptx Pillow requests")
    sys.exit(1)


# ─────────────────────────────────────────────
# 景点数据（与 show_city.py / generate_travel_diary.py 保持同步）
# ─────────────────────────────────────────────
ATTRACTION_DATA = {
    "上海": {
        "外滩": {
            "name_en": "The Bund, Shanghai",
            "description": "外滩是上海最具代表性的历史地标，沿黄浦江西岸绵延约1.6公里，两侧矗立着52栋风格各异的万国建筑群。对岸的陆家嘴摩天楼群与这片旧日繁华遥遥相望，构成了魔都最震撼的天际线。",
            "tips": "建议傍晚至夜间游览，华灯璀璨。可乘坐黄浦江游船从水面欣赏双岸夜景。",
            "best_time": "傍晚至夜间",
            "ticket": "免费开放",
            "wikimedia_file": "Bund_at_night_Shanghai.jpg",
            "image_filename": "waitan.jpg"
        },
        "东方明珠广播电视塔": {
            "name_en": "Oriental Pearl Tower",
            "description": "高耸入云的东方明珠塔是上海浦东的精神图腾，高468米，设计灵感源自唐代诗人白居易的名句《大珠小珠落玉盘》。共有15个观光层，站在玻璃地板观景台俯瞰整个上海，脚下是滚滚黄浦江，远望是无边的都市画卷。",
            "tips": "购买联票同时参观塔内《上海城市历史发展陈列馆》，了解上海百年演变历史。",
            "best_time": "下午3点后（可先拍日景再等夜景）",
            "ticket": "约人民币180元（含全程观光）",
            "wikimedia_file": "Oriental_Pearl_Tower_2012.jpg",
            "image_filename": "oriental_pearl.jpg"
        },
        "豫园": {
            "name_en": "Yu Garden",
            "description": "豫园是上海保存最完好的江南古典园林，始建于1559年明代，由官员潘允端为孝敬父亲而打造，《豫》字取《豫悦老亲》之意。园内湖光山石、亭台楼阁曲折相连，龙墙蜿蜒起伏，四百余年古银杏苍劲挺拔。",
            "tips": "游园后别错过隔壁豫园商城的城隍庙小吃，南翔小笼包是必试经典。",
            "best_time": "工作日上午开门时（避开周末高峰）",
            "ticket": "约人民币40元",
            "wikimedia_file": "Yuyuan_Garden.jpg",
            "image_filename": "yuyuan.jpg"
        }
    },
    "商丘": {
        "商丘古城": {
            "name_en": "Historic City of Shangqiu",
            "description": "商丘古城，又称归德府城，始建于明正德六年（1511年），是中国保存最完整的明清古城之一。古城呈龟背形，城墙高大巍峨，外有护城河环绕。此地叠压着商、周、汉、宋、明五个朝代的古城遗址，被誉为「城摞城」的千年奇观。",
            "tips": "可登上城墙步道漫步，傍晚灯光亮起时最为壮观。",
            "best_time": "春秋两季；傍晚城墙灯光最美",
            "ticket": "景区免费，部分建筑有门票",
            "wikimedia_file": "20220726_Gongchen_Gate,_Historic_City_of_Shangqiu.jpg",
            "image_filename": "shangqiu_old_city.jpg"
        },
        "阏伯台（火神台）": {
            "name_en": "Ebo Tai (Fire God Platform)",
            "description": "阏伯台，又名火神台，是中国现存最早的天文台遗址，距今已有四千余年历史。高台呈圆锥形，高35米，拾级而上望尽豫东平原。每逢农历正月初七祭祀大典，万人赶庙会，是商丘最重要的民间节俗。",
            "tips": "农历正月初七是火神庙会，人气最旺；清晨登台看日出别有意境。",
            "best_time": "农历正月初七庙会期间；清晨看日出",
            "ticket": "约人民币30元",
            "wikimedia_file": "WV_Shangqiu_Banner.jpg",
            "image_filename": "ebo_tai.jpg"
        },
        "芒砀山汉梁王陵": {
            "name_en": "Liang State Tombs at Mount Mangdang",
            "description": "芒砀山是豫东平原唯一山地，因刘邦在此斩白蛇起义而名垂青史。山中密布西汉梁国王室陵墓群，梁孝王陵号称「天下石室第一陵」。陵中壁画《四神云气图》被誉为「敦煌前之敦煌」，是国宝级文物。",
            "tips": "王陵内部全年恒温，夏天游览凉爽；建议穿舒适步行鞋，陵内通道较低。",
            "best_time": "春季（3-5月）或秋季（9-11月）",
            "ticket": "约人民币120元（含景区交通）",
            "wikimedia_file": "Liang_State_Tombs,_Mangdang_Mountain_1.jpg",
            "image_filename": "mangdang_mountain.jpg"
        }
    },
    "北京": {
        "故宫（紫禁城）": {
            "name_en": "The Forbidden City (Palace Museum)",
            "description": "故宫是中国明清两朝的皇家宫殿，建于1406年，历经24位皇帝居于此处达505年。占地72公顷，现存建筑980座，是世界上规模最大、保存最完整的古代木结构建筑群。中心轴线纵贯全宫，金瓦红墙层叠起伏，是中国封建皇权制度的最高展现。",
            "tips": "提前在官方小程序预约门票，现场不售票。建议上午开门时即到避开人流高峰。",
            "best_time": "春秋两季；工作日上午开门时人流较少",
            "ticket": "淡季人民币60元，旺季人民币80元",
            "wikimedia_file": "The_Forbidden_City_-_View_from_Coal_Hill.jpg",
            "image_filename": "forbidden_city.jpg"
        },
        "颐和园": {
            "name_en": "Summer Palace (Yiheyuan)",
            "description": "颐和园是中国保存最完整的皇家园林，1998年列入世界文化遗产。园内以万寿山和昆明湖为核心，佛香阁屹立山中、七百二十八米长廊延展湖临，近代历史与自然山水入画一般合谐共存。",
            "tips": "建议乘坐游船游昆明湖。从东宫门走进，沿长廊漫步至石舫，景色绝佳。",
            "best_time": "春季（三五月桃花盛开）或秋季（10月红叶最美）",
            "ticket": "入园票人民币30元，联票人民币60元（含内部景点）",
            "wikimedia_file": "Kina_Summerpalace_117.jpg",
            "image_filename": "summer_palace.jpg"
        },
        "天坛": {
            "name_en": "Temple of Heaven (Tiantan)",
            "description": "天坛是明清两朝帝王祭天、祈谷的神圣场所，建于1420年，与北京故宫同年建成。祈年殿三层蓝色琉璃屋顶呈圆形，不使一款铁钉纯用榫卯结构，是中国古代建筑的巅峰之作。圆丘坛天心石上发出的声音层层反射，举世奇绝。",
            "tips": "公园早晨有各种民间技艺表演，早起进园欣赏。游览顺序：圆丘坛→皇穹宇→祈年殿。",
            "best_time": "春天（天气晴朗）或秋天（蓝天红叶映衬祈年殿）最美",
            "ticket": "入园票人民币15元，联票人民币34元（含内部建筑）",
            "wikimedia_file": "TempleofHeaven-HallofPrayer.jpg",
            "image_filename": "temple_of_heaven.jpg"
        }
    },
    "南京": {
        "中山陵": {
            "name_en": "Sun Yat-sen Mausoleum",
            "description": "中山陵是中国革命先行者孙中山先生的陵墓，建于1926—1929年，依钟山南坡而建。392级石阶象征着艰辛的革命历程，登至祭堂可俯瞰南京城全貌。气势雄伟、庄严肃穆，是中国近代史上最重要的纪念建筑之一，也是南京必游地标。",
            "tips": "建议徒步攀登392级台阶感受革命精神；清晨或傍晚游人较少，光线最美。可与明孝陵合并游览。",
            "best_time": "春季（梅花盛开3-4月）或秋季（红叶10-11月）",
            "ticket": "免费开放",
            "wikimedia_file": "Sun_yatse_mausoleum.jpg",
            "image_filename": "zhongshaling.jpg"
        },
        "明孝陵": {
            "name_en": "Ming Xiaoling Mausoleum",
            "description": "明孝陵是明朝开国皇帝朱元璋及皇后马氏的合葬陵墓，建于1381年，是中国现存最大的明代皇陵。陵前神道两侧排列着12对巨型石兽与6对石人，气势恢宏。2003年列入联合国教科文组织世界遗产名录，是南京最重要的历史遗迹之一。",
            "tips": "从下马坊入口进，沿神道走至宝城；与中山陵相邻，适合合并游览，全程步行约需2—3小时。",
            "best_time": "春季赏梅（2-3月），秋季红叶（10-11月）最佳",
            "ticket": "约人民币70元",
            "wikimedia_file": "Ming_Xiaoling_Mausoleum_Spirit_Way.jpg",
            "image_filename": "ming_xiaoling.jpg"
        },
        "夫子庙秦淮风光带": {
            "name_en": "Confucius Temple & Qinhuai River Scenic Area",
            "description": "夫子庙，古称金陵文庙，始建于1034年北宋，是供奉和祭祀孔子的庙宇。庙前的秦淮河是六朝以来的风流繁华之地，河畔画舫穿梭，两岸酒肆林立。夜晚华灯璀璨倒映水中，是南京最具烟火气的历史街区，灯火辉煌令人流连忘返。",
            "tips": "傍晚乘坐秦淮画舫游览夜景，别错过桂花鸭、状元豆等本地小吃；元宵节灯会时期最为壮观。",
            "best_time": "傍晚至夜间（灯光最美）；元宵节灯会尤为壮观",
            "ticket": "景区免费，游船约人民币100元",
            "wikimedia_file": "Fuzimiao_nanjing.JPG",
            "image_filename": "fuzimiao.jpg"
        }
    },
    "洛阳": {
        "龙门石窟": {
            "name_en": "Longmen Grottoes",
            "description": "龙门石窟是中国三大石刻艺术宝库之一，2000年列入联合国教科文组织世界文化遗产。两山崖壁上密布着2300余个石龛、10万余尊佛像，最高的卢舍那大佛高达17.14米，面容慈悲安详，被誉为'东方蒙娜丽莎'。伊河两岸青山如黛，石窟与山水相映，历经1500年的时光沉淀，至今仍震撼人心。",
            "tips": "建议游览西山石窟为主，重点参观奉先寺卢舍那大佛；夜游龙门灯光秀极为壮观，白天游览适合拍照。",
            "best_time": "春季（4月牡丹盛开）或秋季（10月红叶金黄）",
            "ticket": "约人民币90元",
            "wikimedia_file": "Longmen Grottoes, Luoyang, Henan.jpg",
            "image_filename": "longmen_grottoes.jpg"
        },
        "白马寺": {
            "name_en": "White Horse Temple",
            "description": "白马寺始建于东汉永平十一年（公元68年），是中国第一座官方修建的佛教寺院，素有'中国第一古刹'之称。相传汉明帝遣使西域求法，两位印度高僧以白马驮载佛经来华，皇帝特建此寺以纪念。寺内现存建筑多为元明时期遗构，古朴庄严，近2000年香火不断。",
            "tips": "寺内有来自印度、缅甸、泰国的佛殿，风格迥异；清晨梵钟响起时氛围最为庄严肃穆。",
            "best_time": "清晨（晨钟暮鼓最具禅意）",
            "ticket": "约人民币50元",
            "wikimedia_file": "White Horse Temple 2.jpg",
            "image_filename": "white_horse_temple.jpg"
        },
        "关林": {
            "name_en": "Guanlin Temple",
            "description": "关林是埋葬三国名将关羽首级之处，也是中国唯一以'林'命名的武将冢，始建于汉代，明万历年间扩建为现今规模。殿宇宏伟、古柏参天，石狮石碑成排列队，气势肃穆。关羽被历代帝王封为'武圣'，关林是全球华人祭祀关公的重要圣地，香火旺盛。",
            "tips": "门前立有清代石碑72通，碑林极具历史价值；正月十三关公诞辰庙会人气极旺，是感受民俗的好时机。",
            "best_time": "全年皆宜；正月十三庙会热闹非凡",
            "ticket": "约人民币40元",
            "wikimedia_file": "Guanlin Temple, Luoyang - September 2011 (6154312540).jpg",
            "image_filename": "guanlin.jpg"
        }
    },
    "青岛": {
        "栈桥": {
            "name_en": "Zhanqiao Pier",
            "description": "栈桥是青岛最具代表性的历史地标，建于1891年，是青岛最早的军事专用人工码头。栈桥伸入胶州湾约440米，桥端建有回澜阁，碧瓦飞檐倒映海中；海鸥翻飞、涛声阵阵，是感受青岛海滨城市气质的最佳起点，也是每一位来青岛游客必打卡的地方。",
            "tips": "清晨或傍晚人少，光线最佳；可在附近的第六海水浴场散步，欣赏红瓦绿树的老城风光。",
            "best_time": "全年皆宜；夏季（6-8月）海景最美",
            "ticket": "免费开放",
            "wikimedia_file": "Qingdao Zhan Qiao.jpg",
            "image_filename": "zhanqiao.jpg"
        },
        "八大关景区": {
            "name_en": "Badaguan Scenic Area",
            "description": "八大关是青岛著名的老城别墅区，因区内有八条以长城各关口命名的道路而得名。这里汇聚了德、英、法、俄、日等20余国风格的各式别墅300余栋，被誉为'万国建筑博览会'。四季花木轮换，春赏樱花、夏看蔷薇、秋观银杏、冬品枫叶，是青岛最诗意的漫步之地。",
            "tips": "适合漫步游览，全程约1-2小时；建议春季（4月）来看道路两侧的樱花盛况最为壮观。",
            "best_time": "春季（4月樱花）或秋季（10月银杏）最美",
            "ticket": "景区免费开放",
            "wikimedia_file": "Badaguan, Qingdao.jpg",
            "image_filename": "badaguan.jpg"
        },
        "崂山": {
            "name_en": "Mount Laoshan",
            "description": "崂山是中国海拔最高的海岸山脉，主峰巨峰海拔1132.7米，素有'海上名山第一'之称。山中泉水清冽甘甜，传说是用来酿造青岛啤酒的水源，道教宫观太清宫在此延续千年香火。登顶俯瞰波澜壮阔的黄海，山海相连的壮景令人叹为观止。",
            "tips": "建议乘坐景区观光车减少体力消耗；太清宫是必游之处，山顶观日出极为震撼，需提前出发。",
            "best_time": "春秋两季（避开夏季高温和冬季山顶寒冷）",
            "ticket": "约人民币130元",
            "wikimedia_file": "Laoshan, Qingdao, Shandong, China - panoramio (32).jpg",
            "image_filename": "laoshan.jpg"
        }
    },
    "广州": {
        "广州塔": {
            "name_en": "Canton Tower",
            "description": "广州塔是广州市的地标性建筑，总高600米（含天线桅杆），是中国第一高塔、世界第四高塔。塔身楼层为454米，外形独特，以两个椭圆形斜交网格结构呈'扭腰'造型，被昵称为'小蛮腰'。夜晚灯光秀璀璨夺目，映照珠江，是广州最浪漫的夜景地标。",
            "tips": "建议傍晚前往，可先看日落后欣赏夜景灯光秀；珠江新城一侧观望效果极佳，游船游览珠江时视角更佳。",
            "best_time": "傍晚至夜间（灯光秀最精彩）",
            "ticket": "约人民币150-200元（根据观光层不同而定）",
            "wikimedia_file": "Canton_Tower_2017-10-02.jpg",
            "image_filename": "canton_tower.jpg"
        },
        "陈家祠": {
            "name_en": "Chen Clan Ancestral Hall",
            "description": "陈家祠（陈氏书院）建于1894年清代，是广东省规模最大、保存最完整的传统岭南建筑群，也是广州民间工艺博物馆。建筑内外密布木雕、砖雕、石雕、灰塑、陶塑及铁铸工艺，七十二间厅堂庄严气派，是岭南建筑艺术的集大成者，被誉为'岭南建筑艺术明珠'。",
            "tips": "馆内集中展示广绣、广彩、广雕等非遗工艺；周末常有粤剧演出，建议提前查询演出时间。",
            "best_time": "全年皆宜；工作日上午人流较少",
            "ticket": "约人民币10元",
            "wikimedia_file": "Chen_Clan_Ancestral_Hall_2025.06_03.jpg",
            "image_filename": "chen_clan_hall.jpg"
        },
        "中山纪念堂": {
            "name_en": "Sun Yat-sen Memorial Hall",
            "description": "广州中山纪念堂是为纪念中国民主革命先驱孙中山先生而建，落成于1931年，是中国著名的宏伟建筑之一。整座建筑呈八角形，采用中国传统宫殿式风格，高47米，跨度71米，全场无一根支柱。主厅可容纳5000余人，宏伟壮观，周边园林绿树葱茏，与越秀山连为一体。",
            "tips": "纪念堂周边的越秀公园可步行前往，五羊雕塑就在附近；建议上午进园，城墙遗址和木兰石须一看。",
            "best_time": "秋冬（10月至次年2月天气舒适）",
            "ticket": "免费开放（节假日可能有活动需购票）",
            "wikimedia_file": "Sun_Yat-sen_Memorial_Hall_Guangzhou.jpg",
            "image_filename": "sun_yatsen_hall.jpg"
        }
    }
}

# ─── 配色方案 ───
COLOR_BG        = RGBColor(0x1A, 0x1A, 0x2E)   # 深蓝背景
COLOR_ACCENT    = RGBColor(0xE9, 0x4F, 0x37)   # 橙红强调
COLOR_GOLD      = RGBColor(0xF5, 0xC9, 0x18)   # 金色标题
COLOR_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_LIGHT     = RGBColor(0xCC, 0xCC, 0xDD)
COLOR_CARD      = RGBColor(0x16, 0x21, 0x3E)   # 卡片背景


def get_wikimedia_direct_url(filename: str) -> str:
    try:
        api_url = "https://commons.wikimedia.org/w/api.php"
        params = {
            "action": "query",
            "titles": f"File:{filename}",
            "prop": "imageinfo",
            "iiprop": "url",
            "iiurlwidth": 1200,
            "format": "json"
        }
        headers = {"User-Agent": "TravelPPTBot/1.0 (educational project)"}
        resp = requests.get(api_url, params=params, headers=headers, timeout=15)
        resp.raise_for_status()
        data = resp.json()
        pages = data.get("query", {}).get("pages", {})
        for page in pages.values():
            info = page.get("imageinfo", [])
            if info:
                return info[0].get("thumburl") or info[0].get("url", "")
        return ""
    except Exception:
        return ""


def download_image(wikimedia_file: str, save_path: Path) -> bool:
    if save_path.exists() and save_path.stat().st_size > 1000:
        return True
    url = get_wikimedia_direct_url(wikimedia_file)
    if not url:
        return False
    try:
        headers = {"User-Agent": "TravelPPTBot/1.0"}
        resp = requests.get(url, headers=headers, timeout=30, stream=True)
        resp.raise_for_status()
        with open(save_path, "wb") as f:
            for chunk in resp.iter_content(chunk_size=8192):
                f.write(chunk)
        Image.open(save_path).verify()
        return True
    except Exception:
        if save_path.exists():
            save_path.unlink()
        return False


def create_folder_structure(base_dir: Path, city: str, attractions: list) -> dict:
    city_dir = base_dir / "旅行日记" / city
    images_dir = city_dir / "景点图片"
    ppt_dir = city_dir / "旅行PPT"
    for d in [city_dir, images_dir, ppt_dir]:
        d.mkdir(parents=True, exist_ok=True)
        print(f"  ✓ 已创建文件夹：{d}")
    attraction_dirs = {}
    for name in attractions:
        ad = images_dir / name
        ad.mkdir(exist_ok=True)
        attraction_dirs[name] = ad
        print(f"  ✓ 已创建景点文件夹：{ad}")
    return {"city_dir": city_dir, "images_dir": images_dir,
            "ppt_dir": ppt_dir, "attraction_dirs": attraction_dirs}


def add_bg(slide, prs):
    """为幻灯片填充深色背景"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BG


def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=None, align=PP_ALIGN.LEFT,
                word_wrap=True):
    from pptx.util import Emu
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color or COLOR_WHITE
    run.font.name = "Microsoft YaHei"
    return txBox


def add_rect(slide, left, top, width, height, fill_color, alpha=None):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def generate_ppt(city: str, base_dir: Path) -> Path:
    attractions_info = ATTRACTION_DATA.get(city, {})
    if not attractions_info:
        print(f"  ✗ 暂无「{city}」的景点数据，请先添加到 ATTRACTION_DATA")
        sys.exit(1)

    attraction_names = list(attractions_info.keys())
    paths = create_folder_structure(base_dir, city, attraction_names)

    # 下载图片
    print("\n  下载景点图片...")
    img_paths = {}
    for name, info in attractions_info.items():
        img_path = paths["attraction_dirs"][name] / info["image_filename"]
        ok = download_image(info["wikimedia_file"], img_path)
        if ok:
            print(f"  ✓ {info['image_filename']}")
        else:
            print(f"  ✗ 下载失败：{info['image_filename']}")
        img_paths[name] = img_path if ok else None

    # 创建 PPT（16:9 宽屏）
    print("\n  生成 PPT...")
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    W = prs.slide_width
    H = prs.slide_height

    blank_layout = prs.slide_layouts[6]  # 完全空白

    # ── 封面页 ──────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide, prs)

    # 装饰色条（左侧竖条）
    add_rect(slide, 0, 0, Inches(0.15), H, COLOR_ACCENT)

    # 城市名大字
    add_textbox(slide, city,
                Inches(0.5), Inches(1.5), Inches(12), Inches(2.0),
                font_size=80, bold=True, color=COLOR_GOLD,
                align=PP_ALIGN.CENTER)

    # 副标题
    add_textbox(slide, "旅  行  日  记",
                Inches(0.5), Inches(3.6), Inches(12), Inches(0.8),
                font_size=28, bold=False, color=COLOR_WHITE,
                align=PP_ALIGN.CENTER)

    # 日期
    date_str = datetime.now().strftime("%Y · %m · %d")
    add_textbox(slide, date_str,
                Inches(0.5), Inches(4.5), Inches(12), Inches(0.6),
                font_size=18, color=COLOR_LIGHT,
                align=PP_ALIGN.CENTER)

    # 底部装饰线
    add_rect(slide, Inches(3), Inches(5.5), Inches(7), Inches(0.04), COLOR_ACCENT)

    add_textbox(slide, "探索世界的美好  用镜头与文字记录每一段旅程",
                Inches(1), Inches(5.8), Inches(11), Inches(0.5),
                font_size=14, color=COLOR_LIGHT,
                align=PP_ALIGN.CENTER)

    # ── 目录页 ──────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide, prs)
    add_rect(slide, 0, 0, Inches(0.15), H, COLOR_GOLD)

    add_textbox(slide, "本次游览景点",
                Inches(0.4), Inches(0.3), Inches(10), Inches(0.8),
                font_size=32, bold=True, color=COLOR_GOLD)
    add_rect(slide, Inches(0.4), Inches(1.2), Inches(6), Inches(0.04), COLOR_ACCENT)

    for i, name in enumerate(attraction_names):
        info = attractions_info[name]
        y = Inches(1.5) + i * Inches(1.5)
        # 序号圆圈背景
        add_rect(slide, Inches(0.4), y, Inches(0.6), Inches(0.6), COLOR_ACCENT)
        add_textbox(slide, str(i + 1),
                    Inches(0.4), y, Inches(0.6), Inches(0.6),
                    font_size=20, bold=True, color=COLOR_WHITE,
                    align=PP_ALIGN.CENTER)
        add_textbox(slide, name,
                    Inches(1.2), y, Inches(5), Inches(0.55),
                    font_size=24, bold=True, color=COLOR_WHITE)
        add_textbox(slide, info["name_en"],
                    Inches(1.2), y + Inches(0.5), Inches(8), Inches(0.45),
                    font_size=14, color=COLOR_LIGHT)

    # ── 每个景点：2张幻灯片 ──────────────────────────

    for idx, (name, info) in enumerate(attractions_info.items(), 1):
        img_path = img_paths.get(name)

        # 景点封面幻灯片（全屏图片 + 标题）
        slide = prs.slides.add_slide(blank_layout)
        add_bg(slide, prs)

        if img_path and img_path.exists():
            try:
                # 图片铺满右侧 2/3
                pic = slide.shapes.add_picture(
                    str(img_path),
                    Inches(4.5), Inches(0),
                    Inches(8.83), H
                )
                # 左侧渐变遮罩（模拟）
                add_rect(slide, Inches(4.0), 0, Inches(1.5), H, COLOR_BG)
            except Exception as e:
                print(f"  ⚠ 图片插入失败：{e}")

        # 左侧：序号 + 标题
        add_rect(slide, 0, 0, Inches(0.15), H, COLOR_ACCENT)
        add_textbox(slide, f"0{idx}",
                    Inches(0.4), Inches(0.5), Inches(3.8), Inches(1.2),
                    font_size=64, bold=True, color=COLOR_ACCENT,
                    align=PP_ALIGN.LEFT)
        add_textbox(slide, name,
                    Inches(0.4), Inches(1.8), Inches(4.0), Inches(1.0),
                    font_size=36, bold=True, color=COLOR_GOLD)
        add_textbox(slide, info["name_en"],
                    Inches(0.4), Inches(2.9), Inches(4.0), Inches(0.55),
                    font_size=14, color=COLOR_LIGHT)
        add_rect(slide, Inches(0.4), Inches(3.55), Inches(3.5), Inches(0.04), COLOR_ACCENT)

        # 景点信息简条
        add_textbox(slide, f"最佳时间：{info['best_time']}",
                    Inches(0.4), Inches(3.8), Inches(4.0), Inches(0.45),
                    font_size=13, color=COLOR_LIGHT)
        add_textbox(slide, f"门票：{info['ticket']}",
                    Inches(0.4), Inches(4.3), Inches(4.0), Inches(0.45),
                    font_size=13, color=COLOR_LIGHT)

        # 景点详情幻灯片（介绍 + 贴士）
        slide2 = prs.slides.add_slide(blank_layout)
        add_bg(slide2, prs)
        add_rect(slide2, 0, 0, Inches(0.15), H, COLOR_GOLD)

        # 标题
        add_textbox(slide2, name,
                    Inches(0.4), Inches(0.3), Inches(9), Inches(0.75),
                    font_size=30, bold=True, color=COLOR_GOLD)
        add_rect(slide2, Inches(0.4), Inches(1.15), Inches(8), Inches(0.04), COLOR_ACCENT)

        # 插图（右侧小图）
        if img_path and img_path.exists():
            try:
                slide2.shapes.add_picture(
                    str(img_path),
                    Inches(8.8), Inches(0.8),
                    Inches(4.2), Inches(3.5)
                )
            except Exception:
                pass

        # 介绍文字
        add_textbox(slide2, "景点介绍",
                    Inches(0.4), Inches(1.35), Inches(8), Inches(0.5),
                    font_size=16, bold=True, color=COLOR_ACCENT)
        add_textbox(slide2, info["description"],
                    Inches(0.4), Inches(1.9), Inches(8.2), Inches(2.2),
                    font_size=15, color=COLOR_WHITE, word_wrap=True)

        # 贴士卡片
        add_rect(slide2, Inches(0.4), Inches(4.3), Inches(8.2), Inches(1.4), COLOR_CARD)
        add_textbox(slide2, "旅行贴士",
                    Inches(0.55), Inches(4.45), Inches(7.8), Inches(0.45),
                    font_size=14, bold=True, color=COLOR_GOLD)
        add_textbox(slide2, info["tips"],
                    Inches(0.55), Inches(4.9), Inches(7.8), Inches(0.75),
                    font_size=13, color=COLOR_LIGHT, word_wrap=True)

    # ── 结尾页 ──────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide, prs)
    add_rect(slide, 0, 0, Inches(0.15), H, COLOR_ACCENT)
    add_rect(slide, 0, H - Inches(0.15), W, Inches(0.15), COLOR_ACCENT)

    add_textbox(slide, "旅途未完，精彩继续",
                Inches(0.5), Inches(2.0), Inches(12.3), Inches(1.2),
                font_size=42, bold=True, color=COLOR_GOLD,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, f"感谢探索  {city}  的每一处美好",
                Inches(0.5), Inches(3.4), Inches(12.3), Inches(0.7),
                font_size=22, color=COLOR_WHITE,
                align=PP_ALIGN.CENTER)
    add_rect(slide, Inches(4), Inches(4.3), Inches(5), Inches(0.04), COLOR_ACCENT)
    add_textbox(slide, "由「虚拟旅行」技能自动生成",
                Inches(0.5), Inches(4.6), Inches(12.3), Inches(0.5),
                font_size=14, color=COLOR_LIGHT,
                align=PP_ALIGN.CENTER)

    # 保存
    date_str = datetime.now().strftime("%Y%m%d")
    ppt_filename = f"{city}_旅行日记_{date_str}.pptx"
    ppt_path = paths["ppt_dir"] / ppt_filename
    prs.save(str(ppt_path))
    print(f"\n  ✅ PPT 已生成：{ppt_path}")
    return ppt_path


def main():
    parser = argparse.ArgumentParser(description="旅行日记 PPT 生成器")
    parser.add_argument("--city", required=True, help="目标城市名称")
    parser.add_argument("--output-dir", default=".", help="输出根目录")
    args = parser.parse_args()

    base_dir = Path(args.output_dir).resolve()
    city = args.city

    print(f"\n  生成「{city}」旅行日记 PPT...\n")
    print("  创建文件夹结构...")
    ppt_path = generate_ppt(city, base_dir)
    print(f"\n  完成！文件位于：{ppt_path}")
    return str(ppt_path)


if __name__ == "__main__":
    main()
